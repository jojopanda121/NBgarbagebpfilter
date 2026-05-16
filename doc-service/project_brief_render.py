"""
project_brief_render.py (PE/VC 版)
───────────────────────────────────
项目简报 3 页 deck · 一级市场视角 · 确定性渲染器

用法:
    python project_brief_render.py content.json output.pptx

设计原则: 与 investment_snapshot_render.py 同源
    · 模块 / 版式 / 坐标 / 颜色 / 字号锁死.
    · 颜色 / 字体 / 字号一律从 brand_tokens 导入, 与网页 :root 同源.
    · Agent 只产符合 content_schema.json (PE/VC 版) 的 JSON.

3 页结构:
    P1 封面    company_full_name + tagline + metadata 3 chips + dealroom_meta 3 chips
    P2 概况+亮 overview + highlights × 4 (2x2)
    P3 三件套  team × 2-3 + financials_compact 3x4 + valuation_view + risks
"""

import json
import sys

from pptx import Presentation
from pptx.util import Inches

from brand_tokens import (
    COLOR, FONT_CN_SERIF, FONT_CN_SANS, FONT_EN, SIZE,
    RULE_PT, HAIRLINE_PT,
    add_rect, set_run, add_text, add_para,
)


# 画布: A4 横版 (与 investment_snapshot 一致)
SLIDE_W = Inches(11.69)
SLIDE_H = Inches(8.27)

PAGE_TOTAL = 3


def _add_page_footer(slide, page_num, page_total, company):
    geom = {"x": 0.82, "y": 7.85, "w": 10.05, "h": 0.30}
    tf = add_text(slide, geom, align="l", anchor="b", margin=0.05)
    p = add_para(tf, "l")
    set_run(p.add_run(),
            f"{company}　·　项目简报　·　{page_num} / {page_total}",
            size=SIZE["footer"], color=COLOR["mid"])


def _new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render_page_title(slide, text, y=0.42):
    tf = add_text(slide, {"x": 0.82, "y": y, "w": 10.05, "h": 0.40},
                  align="l", margin=0)
    set_run(add_para(tf, "l").add_run(), text,
            size=SIZE["page_title"], bold=True, color=COLOR["navy"],
            font_cn=FONT_CN_SERIF)
    add_rect(slide, {"x": 0.82, "y": y + 0.44, "w": 10.05, "h": 0.05},
             fill=COLOR["accent"])


def _section_label(slide, geom, text):
    """段标签 (与 investment_snapshot 同款): 浅蓝灰胶囊 + 品牌蓝文字"""
    add_rect(slide, geom, fill=COLOR["bg3"])
    tf = add_text(slide, geom, align="c", anchor="m", margin=0)
    set_run(add_para(tf, "c").add_run(), text,
            size=SIZE["section"], bold=True, color=COLOR["accent"])


# ═══════════════════════════════════════════════════════════════════
# P1: 封面 — 公司名 + tagline 横幅 + metadata × 3 + dealroom_meta × 3
# ═══════════════════════════════════════════════════════════════════

def _render_cover(prs, content):
    slide = _new_slide(prs)
    # 顶部标题
    tf = add_text(slide, {"x": 0.82, "y": 1.40, "w": 10.05, "h": 0.85},
                  align="l", margin=0)
    p = add_para(tf, "l")
    set_run(p.add_run(), content["company_full_name"],
            size=SIZE["cover_title"], bold=True, color=COLOR["navy"],
            font_cn=FONT_CN_SERIF)
    # 标题下品牌蓝细线
    add_rect(slide, {"x": 0.82, "y": 2.38, "w": 10.05, "h": 0.06},
             fill=COLOR["accent"])
    # tagline navy 横幅
    add_rect(slide, {"x": 0.82, "y": 2.70, "w": 10.05, "h": 0.62},
             fill=COLOR["navy"])
    tf = add_text(slide, {"x": 0.82, "y": 2.70, "w": 10.05, "h": 0.62},
                  align="c", anchor="m", margin=0.1)
    p = add_para(tf, "c")
    set_run(p.add_run(), content["tagline"],
            size=SIZE["cover_tag"], bold=True, color=COLOR["white"])

    # metadata 3 chips (行业 / 阶段 / 地点)
    meta = content["metadata"]
    _draw_chip_row(
        slide, y=3.85,
        items=[("行业", meta["industry"]),
               ("阶段", meta["stage"]),
               ("地点", meta["location"])],
        label_color=COLOR["accent"],
    )
    # dealroom_meta 3 chips (本轮规模 / Pre 估值 / 领投状态) — 一级市场专属
    dm = content["dealroom_meta"]
    _draw_chip_row(
        slide, y=5.05,
        items=[("本轮规模",   dm["round_size"]),
               ("Pre 估值",   dm["pre_valuation"]),
               ("领投状态",   dm["lead_investor_status"])],
        label_color=COLOR["navy"],
    )

    _add_page_footer(slide, 1, PAGE_TOTAL, content["company_full_name"])


def _draw_chip_row(slide, *, y, items, label_color):
    """3 个胶囊横排, 等宽居中. 标签色由调用方决定 (accent 或 navy)."""
    cap_w, cap_h, gap = 3.10, 0.95, 0.18
    total_w = cap_w * 3 + gap * 2
    x0 = (11.69 - total_w) / 2
    for i, (k, v) in enumerate(items):
        x = x0 + i * (cap_w + gap)
        add_rect(slide, {"x": x, "y": y, "w": cap_w, "h": cap_h},
                 fill=COLOR["bg3"],
                 line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
        # 上半: 标签
        tf = add_text(slide, {"x": x, "y": y + 0.10, "w": cap_w, "h": 0.32},
                      align="c", anchor="m", margin=0.04)
        set_run(add_para(tf, "c").add_run(), k,
                size=SIZE["cover_meta"], bold=True, color=label_color)
        # 下半: 值
        tf = add_text(slide, {"x": x, "y": y + 0.42, "w": cap_w, "h": 0.50},
                      align="c", anchor="m", margin=0.04)
        set_run(add_para(tf, "c").add_run(), v,
                size=SIZE["cover_meta"], color=COLOR["ink"])


# ═══════════════════════════════════════════════════════════════════
# P2: 概况 + 投资亮点 2x2
# ═══════════════════════════════════════════════════════════════════

def _render_overview_and_highlights(prs, content):
    slide = _new_slide(prs)
    _render_page_title(slide, "项目概况 & 投资亮点")

    # 概况区
    add_rect(slide, {"x": 0.82, "y": 1.30, "w": 10.05, "h": 1.70},
             fill=COLOR["bg2"],
             line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    _section_label(slide, {"x": 1.20, "y": 1.10, "w": 1.60, "h": 0.30},
                   "项目概况")
    tf = add_text(slide,
                  {"x": 0.96, "y": 1.46, "w": 9.77, "h": 1.50}, margin=0.08)
    set_run(add_para(tf, "l").add_run(), content["overview"],
            size=SIZE["body"], color=COLOR["ink"])

    # 亮点 2x2 网格
    _section_label(slide, {"x": 4.86, "y": 3.18, "w": 1.98, "h": 0.30},
                   "投资亮点")
    grid_x, grid_y = 0.82, 3.40
    cell_w, cell_h, gap_x, gap_y = 4.92, 2.10, 0.21, 0.18
    for i, item in enumerate(content["highlights"]):
        col, row = i % 2, i // 2
        x = grid_x + col * (cell_w + gap_x)
        y = grid_y + row * (cell_h + gap_y)
        add_rect(slide, {"x": x, "y": y, "w": cell_w, "h": cell_h},
                 fill=COLOR["bg2"],
                 line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
        tf = add_text(slide,
                      {"x": x + 0.16, "y": y + 0.10,
                       "w": cell_w - 0.32, "h": cell_h - 0.20},
                      margin=0.04)
        p = add_para(tf, "l", space_after_pt=4)
        set_run(p.add_run(), f"{item['label']}：",
                size=SIZE["label_inline"], bold=True, color=COLOR["accent"])
        set_run(p.add_run(), item["desc"],
                size=SIZE["body"], color=COLOR["ink"])
    _add_page_footer(slide, 2, PAGE_TOTAL, content["company_full_name"])


# ═══════════════════════════════════════════════════════════════════
# P3: PE/VC 决策三件套 — 团队 / 财务紧凑表 / 估值视角
# ═══════════════════════════════════════════════════════════════════

def _render_team_financials_valuation(prs, content):
    slide = _new_slide(prs)
    _render_page_title(slide, "团队 · 财务 · 估值视角")

    # Row 1 团队 (上方 ~2.4")
    _section_label(slide, {"x": 5.00, "y": 1.10, "w": 1.70, "h": 0.30},
                   "核心团队")
    _draw_team_row(slide, content["team"],
                   {"x": 0.82, "y": 1.30, "w": 10.05, "h": 2.20})

    # Row 2 财务紧凑表 (中部 ~1.8")
    _section_label(slide, {"x": 4.86, "y": 3.70, "w": 1.98, "h": 0.30},
                   "财务速览")
    add_rect(slide, {"x": 0.82, "y": 3.90, "w": 10.05, "h": 1.65},
             fill=COLOR["bg2"],
             line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    _draw_financials_compact(
        slide, content["financials_compact"],
        {"x": 0.96, "y": 4.00, "w": 9.77, "h": 1.45},
    )

    # Row 3 估值视角 + 风险 (底部 ~2")
    _section_label(slide, {"x": 4.86, "y": 5.75, "w": 1.98, "h": 0.30},
                   "估值视角")
    add_rect(slide, {"x": 0.82, "y": 5.95, "w": 10.05, "h": 1.78},
             fill=COLOR["bg2"],
             line_color=COLOR["accent"], line_w_pt=HAIRLINE_PT)
    _draw_valuation_view(slide, content["valuation_view"],
                         {"x": 0.96, "y": 6.05, "w": 5.92, "h": 1.58})
    add_rect(slide, {"x": 7.05, "y": 6.05, "w": 0.012, "h": 1.52},
             fill=COLOR["border"])
    _draw_risk_list(slide, content["risks"],
                    {"x": 7.18, "y": 6.05, "w": 3.46, "h": 1.58})

    _add_page_footer(slide, 3, PAGE_TOTAL, content["company_full_name"])


def _draw_team_row(slide, team, geom):
    """2-3 张团队卡, 等宽横排. 卡内: 姓名 (大) / 角色 (品牌蓝) / bio."""
    n = len(team)
    gap = 0.21
    card_w = (geom["w"] - gap * (n - 1)) / n
    for i, m in enumerate(team):
        x = geom["x"] + i * (card_w + gap)
        add_rect(slide, {"x": x, "y": geom["y"], "w": card_w, "h": geom["h"]},
                 fill=COLOR["bg2"],
                 line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
        # 姓名
        tf = add_text(slide,
                      {"x": x + 0.16, "y": geom["y"] + 0.16,
                       "w": card_w - 0.32, "h": 0.42}, margin=0.02)
        set_run(add_para(tf, "l").add_run(), m["name"],
                size=SIZE["section"] + 4, bold=True, color=COLOR["navy"],
                font_cn=FONT_CN_SERIF)
        # 角色
        tf = add_text(slide,
                      {"x": x + 0.16, "y": geom["y"] + 0.62,
                       "w": card_w - 0.32, "h": 0.30}, margin=0.02)
        set_run(add_para(tf, "l").add_run(), m["role"],
                size=SIZE["label_inline"], bold=True, color=COLOR["accent"])
        # 分隔细线
        add_rect(slide,
                 {"x": x + 0.16, "y": geom["y"] + 0.96,
                  "w": card_w - 0.32, "h": 0.012},
                 fill=COLOR["border"])
        # bio
        tf = add_text(slide,
                      {"x": x + 0.16, "y": geom["y"] + 1.04,
                       "w": card_w - 0.32, "h": geom["h"] - 1.16},
                      margin=0.04)
        p = add_para(tf, "l", space_after_pt=2)
        p.line_spacing = 1.25
        set_run(p.add_run(), m["bio"],
                size=SIZE["body"], color=COLOR["ink"])


def _draw_financials_compact(slide, fin, geom):
    """3 列 (期间) × 4 元素 (指标名 + 3 期值) 紧凑表. 顶下双品牌蓝细线."""
    cols = fin["columns"]   # 3 期
    rows = fin["rows"]      # 3 行
    name_w = geom["w"] * 0.30
    cell_w = (geom["w"] - name_w) / len(cols)
    n_rows = len(rows) + 1  # 表头
    row_h = geom["h"] / n_rows
    y0 = geom["y"]

    # 顶线 (品牌蓝)
    add_rect(slide,
             {"x": geom["x"], "y": y0, "w": geom["w"], "h": 0.02},
             fill=COLOR["accent"])
    # 表头
    headers = ["指标 (期间)"] + cols
    col_x = [geom["x"]] + [geom["x"] + name_w + i * cell_w
                            for i in range(len(cols))]
    col_w_list = [name_w] + [cell_w] * len(cols)
    for i, h in enumerate(headers):
        cell = add_text(slide,
                        {"x": col_x[i], "y": y0 + 0.02,
                         "w": col_w_list[i], "h": row_h},
                        align=("l" if i == 0 else "c"), anchor="m", margin=0.04)
        set_run(add_para(cell, "l" if i == 0 else "c").add_run(), h,
                size=SIZE["table"], bold=True, color=COLOR["mid"])
    # 表头下线
    yh = y0 + 0.02 + row_h
    add_rect(slide,
             {"x": geom["x"], "y": yh, "w": geom["w"], "h": 0.02},
             fill=COLOR["accent"])
    # 数据行
    for ri, row in enumerate(rows):
        ry = yh + 0.02 + ri * row_h
        for ci, val in enumerate(row):
            cell = add_text(slide,
                            {"x": col_x[ci], "y": ry,
                             "w": col_w_list[ci], "h": row_h},
                            align=("l" if ci == 0 else "c"),
                            anchor="m", margin=0.04)
            bold = (ci == 0)
            set_run(add_para(cell, "l" if ci == 0 else "c").add_run(), val,
                    size=SIZE["table"], bold=bold, color=COLOR["ink"])


def _draw_valuation_view(slide, vv, geom):
    """3 字段堆叠: comp_anchor / recommended_range / rationale.

    recommended_range 是核心数字, 用 navy 大字突出."""
    # 1. comp_anchor 行 (mid 色小字)
    tf = add_text(slide,
                  {"x": geom["x"], "y": geom["y"],
                   "w": geom["w"], "h": 0.34}, margin=0.04)
    p = add_para(tf, "l")
    set_run(p.add_run(), "可比锚： ",
            size=SIZE["label_inline"], bold=True, color=COLOR["accent"])
    set_run(p.add_run(), vv["comp_anchor"],
            size=SIZE["body"], color=COLOR["ink"])

    # 2. recommended_range 行 (navy 加粗大字, 主数据)
    tf = add_text(slide,
                  {"x": geom["x"], "y": geom["y"] + 0.40,
                   "w": geom["w"], "h": 0.46}, margin=0.04)
    p = add_para(tf, "l")
    set_run(p.add_run(), "建议估值区间： ",
            size=SIZE["label_inline"], bold=True, color=COLOR["accent"])
    set_run(p.add_run(), vv["recommended_range"],
            size=SIZE["cover_tag"], bold=True, color=COLOR["navy"],
            font_cn=FONT_CN_SERIF)

    # 3. rationale 行 (正文)
    tf = add_text(slide,
                  {"x": geom["x"], "y": geom["y"] + 0.94,
                   "w": geom["w"], "h": geom["h"] - 0.94}, margin=0.04)
    p = add_para(tf, "l")
    p.line_spacing = 1.30
    set_run(p.add_run(), "推导逻辑： ",
            size=SIZE["label_inline"], bold=True, color=COLOR["accent"])
    set_run(p.add_run(), vv["rationale"],
            size=SIZE["body"], color=COLOR["ink"])


def _draw_risk_list(slide, risks, geom):
    """紧凑风险列表: 只承载 3 条标签 + 缓释/影响句, 防止挤爆 P3."""
    tf = add_text(slide,
                  {"x": geom["x"], "y": geom["y"],
                   "w": geom["w"], "h": 0.24}, margin=0.02)
    set_run(add_para(tf, "l").add_run(), "核心风险",
            size=SIZE["label_inline"], bold=True, color=COLOR["red"])

    row_h = (geom["h"] - 0.28) / 3
    for i, item in enumerate(risks[:3]):
        y = geom["y"] + 0.30 + i * row_h
        tf = add_text(slide,
                      {"x": geom["x"], "y": y,
                       "w": geom["w"], "h": row_h - 0.03},
                      margin=0.02)
        p = add_para(tf, "l", space_after_pt=1)
        p.line_spacing = 1.05
        set_run(p.add_run(), f"{item['label']}：",
                size=SIZE["table"], bold=True, color=COLOR["red"])
        set_run(p.add_run(), item["desc"],
                size=SIZE["table"], color=COLOR["ink"])


def render(content: dict, out_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _render_cover(prs, content)
    _render_overview_and_highlights(prs, content)
    _render_team_financials_valuation(prs, content)
    prs.save(out_path)


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python project_brief_render.py <content.json> <output.pptx>")
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        content = json.load(f)
    render(content, sys.argv[2])
    print(f"✓ rendered → {sys.argv[2]}")
