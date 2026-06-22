#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
通用文档文本提取器（Smart Document Parser）
支持两种模式：
  pptx — 使用 python-pptx 直接提取幻灯片文本与演讲者备注（无 OCR）
  pdf  — 文字优先策略：PyMuPDF 直接提取；仅当 avg_chars_per_page <= 50 时
          回退 RapidOCR（替代 Tesseract，速度更快、无需系统依赖）
  docx — 提取 Word 段落和表格
  doc  — 旧版 Word 二进制文件的 best-effort 文本提取
  xlsx — 提取 Excel 工作表预览
  csv  — 提取 CSV 预览

用法: python extract_doc.py <path_to_file> <mode>
      mode: pptx | pdf | docx | doc | xlsx | csv
输出: 纯文本到 stdout；错误时 stderr JSON + exit 1
"""

import sys
import os
import re
import json
import csv
import zipfile
from xml.etree import ElementTree as ET


IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp")


def extraction_fallback_text(file_path: str, mode: str, reason: str = "") -> str:
    """最后兜底：保证调用方拿到可进入分析流程的文本，而不是解析异常。"""
    filename = os.path.basename(file_path)
    try:
        size = os.path.getsize(file_path)
    except OSError:
        size = 0
    reason_line = f"解析提示：{reason}" if reason else "解析提示：常规文本层为空或不可读取。"
    return "\n".join([
        "【上传材料解析结果】",
        f"文件名：{filename}",
        f"文件格式：{mode}",
        f"文件大小：{size} 字节",
        reason_line,
        "系统已接收该文件，但自动解析只能取得有限文本。请基于文件名、格式和已提取片段谨慎分析，并在结论中标注材料可读性不足；不要声称已读取到文件中未实际提取出的具体事实。",
    ])


def _ocr_image_bytes(image_bytes: bytes) -> str:
    """Best-effort OCR：依赖缺失或识别失败时返回空字符串，不中断上传体验。"""
    try:
        from rapidocr_onnxruntime import RapidOCR
    except Exception as e:
        print(f"警告: OCR 依赖不可用，跳过内嵌图片 OCR: {e}", file=sys.stderr)
        return ""
    try:
        ocr = RapidOCR()
        result, _ = ocr(image_bytes)
        if not result:
            return ""
        return "\n".join(line[1] for line in result if len(line) > 1 and line[1]).strip()
    except Exception as e:
        print(f"警告: 内嵌图片 OCR 失败: {e}", file=sys.stderr)
        return ""


def _extract_ooxml_media_ocr(file_path: str, media_prefix: str, max_images: int = 40):
    """对 PPTX/DOCX 内嵌图片做 OCR，覆盖整页截图型材料。"""
    parts = []
    try:
        with zipfile.ZipFile(file_path) as zf:
            names = [
                n for n in zf.namelist()
                if n.startswith(media_prefix) and n.lower().endswith(IMAGE_EXTS)
            ]
            for idx, name in enumerate(names[:max_images], start=1):
                text = _ocr_image_bytes(zf.read(name))
                if text:
                    parts.append(f"[图片 OCR {idx}: {os.path.basename(name)}]\n{text}")
            if len(names) > max_images:
                parts.append(f"...（内嵌图片超过 {max_images} 张，其余已跳过 OCR）")
    except Exception as e:
        print(f"警告: 读取 OOXML 内嵌图片失败: {e}", file=sys.stderr)
    return parts


def _ooxml_xml_text(file_path: str, wanted_prefixes):
    """从 OOXML XML 文本节点兜底提取内容。"""
    def natural_key(name):
        return [int(x) if x.isdigit() else x for x in re.split(r"(\d+)", name)]

    with zipfile.ZipFile(file_path) as zf:
        names = sorted(
            (n for n in zf.namelist() if n.endswith(".xml") and n.startswith(tuple(wanted_prefixes))),
            key=natural_key,
        )
        for name in names:
            try:
                root = ET.fromstring(zf.read(name))
            except Exception:
                continue
            for elem in root.iter():
                local = elem.tag.rsplit("}", 1)[-1]
                if local not in {"t", "v", "instrText"}:
                    continue
                text = (elem.text or "").strip()
                if text:
                    yield text


# ─────────────────────────────────────────────────────────────
# Mode A: PPTX — python-pptx 直接提取，无需 OCR
# ─────────────────────────────────────────────────────────────

def _normalize_pptx_text(text: str) -> str:
    text = re.sub(r"[ \t\r\f\v]+", " ", (text or "")).strip()
    return text


def _remember_text(seen, text: str):
    text = _normalize_pptx_text(text)
    if text:
        seen.add(text)


def _append_unique(parts, seen, text: str):
    text = _normalize_pptx_text(text)
    if not text:
        return
    if text in seen:
        return
    seen.add(text)
    parts.append(text)


def _iter_pptx_shape_text(shape):
    """递归提取 PPTX shape 文本，覆盖文本框、表格、组合形状。"""
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            line = para.text.strip()
            if line:
                yield line

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = []
            for cell in row.cells:
                cell_text = (cell.text or "").strip().replace("\n", " ")
                cells.append(cell_text)
            if any(cells):
                yield " | ".join(cells)

    # GroupShape 暴露 .shapes；递归进入组合形状/嵌套形状。
    child_shapes = getattr(shape, "shapes", None)
    if child_shapes:
        for child in child_shapes:
            yield from _iter_pptx_shape_text(child)


def _pptx_xml_fallback_text(file_path: str):
    """
    从 PPTX XML 兜底提取 DrawingML 文本。
    python-pptx 不总能展开 SmartArt、图表缓存、部分嵌套对象；这些文本通常仍在
    ppt/slides、ppt/charts、ppt/diagrams 或 notesSlides XML 内。
    """
    wanted_prefixes = (
        "ppt/slides/slide",
        "ppt/notesSlides/notesSlide",
        "ppt/charts/chart",
        "ppt/diagrams/data",
    )

    yield from _ooxml_xml_text(file_path, wanted_prefixes)


def extract_pptx(file_path: str) -> str:
    """遍历所有幻灯片，提取标题/正文/表格/组合形状/备注，并用 XML 兜底。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("请安装 python-pptx: pip install python-pptx")

    prs = Presentation(file_path)
    parts = []
    seen = set()

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts = [f"--- 第 {slide_idx} 页 ---"]
        slide_seen = set()

        # 遍历所有形状，提取文本框、表格、组合形状等内容。
        for shape in slide.shapes:
            for line in _iter_pptx_shape_text(shape):
                _append_unique(slide_parts, slide_seen, line)
                _remember_text(seen, line)

        # 提取演讲者备注（speaker notes）
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                if notes_text:
                    _append_unique(slide_parts, slide_seen, f"[备注] {notes_text}")
                    _remember_text(seen, notes_text)

        parts.append("\n".join(slide_parts))

    # XML 兜底：补齐 SmartArt / 图表缓存 / 其它 python-pptx 未暴露的文本。
    xml_parts = []
    for line in _pptx_xml_fallback_text(file_path):
        normalized = _normalize_pptx_text(line)
        if any(normalized != existing and normalized in existing for existing in seen):
            continue
        _append_unique(xml_parts, seen, line)
    if xml_parts:
        parts.append("[PPTX XML 兜底文本]\n" + "\n".join(xml_parts))

    current_len = len("\n".join(parts).strip())
    if current_len < 200:
        ocr_parts = _extract_ooxml_media_ocr(file_path, "ppt/media/")
        if ocr_parts:
            parts.append("[PPTX 内嵌图片 OCR]\n" + "\n\n".join(ocr_parts))

    return "\n\n".join(parts).strip()


def extract_docx(file_path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        Document = None

    parts = []
    seen = set()

    if Document:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                _append_unique(parts, seen, text)
        for table_idx, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                _append_unique(parts, seen, f"[表格 {table_idx}]\n" + "\n".join(rows))

    # XML 兜底：页眉页脚、文本框、批注、脚注等 python-docx 可能漏掉的文本。
    xml_parts = []
    wanted_prefixes = (
        "word/document",
        "word/header",
        "word/footer",
        "word/footnotes",
        "word/endnotes",
        "word/comments",
    )
    try:
        for text in _ooxml_xml_text(file_path, wanted_prefixes):
            normalized = _normalize_pptx_text(text)
            if any(normalized != existing and normalized in existing for existing in seen):
                continue
            _append_unique(xml_parts, seen, text)
    except Exception as e:
        print(f"警告: DOCX XML 兜底提取失败: {e}", file=sys.stderr)
    if xml_parts:
        parts.append("[DOCX XML 兜底文本]\n" + "\n".join(xml_parts))

    current_len = len("\n".join(parts).strip())
    if current_len < 200:
        ocr_parts = _extract_ooxml_media_ocr(file_path, "word/media/")
        if ocr_parts:
            parts.append("[DOCX 内嵌图片 OCR]\n" + "\n\n".join(ocr_parts))

    return "\n\n".join(parts).strip()


def extract_legacy_doc(file_path: str) -> str:
    """
    旧版 .doc 是 OLE/CFB 二进制格式。没有 LibreOffice/antiword 时无法完整还原版式，
    这里做 best-effort 字符串提取，至少不让上传流程失败。
    """
    data = open(file_path, "rb").read()
    parts = []
    seen = set()

    for encoding in ("utf-16le", "gb18030", "utf-8", "latin1"):
        try:
            decoded = data.decode(encoding, errors="ignore")
        except Exception:
            continue
        # 保留中英文、数字和常见业务符号的连续片段，过滤二进制噪声。
        for match in re.finditer(r"[\u4e00-\u9fffA-Za-z0-9][\u4e00-\u9fffA-Za-z0-9\s，。、“”‘’：；！？（）()《》<>/%.,:;!?+\-_=·&@#￥$]{3,}", decoded):
            text = re.sub(r"\s+", " ", match.group(0)).strip()
            if len(text) < 4 or text in seen:
                continue
            seen.add(text)
            parts.append(text)
            if len(parts) >= 300:
                break
        if len("\n".join(parts)) > 200:
            break

    return "\n".join(parts).strip()


def extract_xlsx(file_path: str) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("请安装 openpyxl: pip install openpyxl")

    wb = load_workbook(file_path, data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        lines = [f"# Sheet: {ws.title}"]
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if not any(v.strip() for v in values):
                continue
            lines.append(" | ".join(values[:20]))
            row_count += 1
            if row_count >= 80:
                lines.append("...（已截断，仅展示前 80 行非空数据）")
                break
        if row_count:
            parts.append("\n".join(lines))
    wb.close()
    return "\n\n".join(parts).strip()


def extract_csv(file_path: str) -> str:
    encodings = ("utf-8-sig", "utf-8", "gb18030")
    last_error = None
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc, newline="") as f:
                reader = csv.reader(f)
                lines = []
                for idx, row in enumerate(reader):
                    if idx >= 120:
                        lines.append("...（已截断，仅展示前 120 行）")
                        break
                    lines.append(" | ".join(row[:30]))
                return "\n".join(lines).strip()
        except UnicodeDecodeError as e:
            last_error = e
    raise RuntimeError(f"CSV 编码识别失败: {last_error}")


# ─────────────────────────────────────────────────────────────
# Mode B: PDF — 文字优先 + RapidOCR 兜底
# ─────────────────────────────────────────────────────────────

def extract_pdf(file_path: str) -> str:
    """
    Step 1: 用 PyMuPDF 直接抽取文字层。
    Step 2: 计算 avg_chars_per_page = total_text_len / num_pages。
    Step 3: avg_chars_per_page > 50  → 快速路径，直接返回文字。
            avg_chars_per_page <= 50 → 慢速路径，对页面图像运行 RapidOCR。
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return extraction_fallback_text(file_path, "pdf", "PDF 解析依赖 PyMuPDF 不可用，未能读取文本层。")

    doc = fitz.open(file_path)

    # 检测加密 PDF
    if doc.is_encrypted:
        if not doc.authenticate(""):
            doc.close()
            raise RuntimeError("PDF 已加密且需要密码，请上传未加密的 PDF 文件")

    num_pages = len(doc)
    if num_pages == 0:
        doc.close()
        raise RuntimeError("PDF 文件为空（0 页），请检查文件是否完整")

    # Step 1: 提取文字层
    parts = []
    for i in range(num_pages):
        try:
            parts.append(doc[i].get_text() or "")
        except Exception as e:
            print(f"警告: 第 {i+1} 页提取失败: {e}", file=sys.stderr)
            parts.append("")

    text = "\n".join(parts).replace("\r\n", "\n").strip()

    # Step 2: 密度检查
    avg_chars_per_page = len(text) / num_pages

    # Step 3: 决策
    if avg_chars_per_page > 50:
        # 快速路径：文字版 PDF，直接返回
        doc.close()
        return re.sub(r"[ \t]+", " ", text)

    # 慢速路径：扫描版 PDF，启动 RapidOCR
    print(
        f"提示: 每页均字符数 {avg_chars_per_page:.1f} <= 50，启动 RapidOCR...",
        file=sys.stderr,
    )

    try:
        import numpy as np
        from rapidocr_onnxruntime import RapidOCR
    except ImportError as e:
        print(f"警告: OCR 依赖不可用，返回 PDF 文字层: {e}", file=sys.stderr)
        if text:
            doc.close()
            return re.sub(r"[ \t]+", " ", text)
        doc.close()
        return extraction_fallback_text(file_path, "pdf", "PDF 文字层为空，且 OCR 依赖不可用。")

    ocr = RapidOCR()
    ocr_parts = []

    for i in range(num_pages):
        try:
            # 渲染页面为像素图（200 dpi）
            pix = doc[i].get_pixmap(dpi=200)
            # 转换为 numpy 数组（RapidOCR 接受 numpy array）
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                pix.height, pix.width, pix.n
            )
            # PyMuPDF 默认 RGB(3) 或 RGBA(4)，RapidOCR 需要 RGB
            if pix.n == 4:
                img = img[:, :, :3]
            result, _ = ocr(img)
            if result:
                page_text = "\n".join(line[1] for line in result)
                ocr_parts.append(page_text)
            else:
                ocr_parts.append("")
        except Exception as e:
            print(f"警告: 第 {i+1} 页 OCR 失败: {e}", file=sys.stderr)
            ocr_parts.append("")

    doc.close()

    ocr_text = "\n".join(ocr_parts).replace("\r\n", "\n").strip()
    # 若文字层存在，只有 OCR 明显更丰富时才替换，避免短文本 PDF 被 OCR 误伤。
    if text and ocr_text:
        final = ocr_text if len(ocr_text) > len(text) * 1.2 else text
    else:
        final = ocr_text or text
    return re.sub(r"[ \t]+", " ", final)


# ─────────────────────────────────────────────────────────────
# 入口
# ─────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 3:
        print("用法: python extract_doc.py <path_to_file> <mode>", file=sys.stderr)
        print("  mode: pptx | pdf | docx | doc | xlsx | csv", file=sys.stderr)
        sys.exit(1)

    file_path = sys.argv[1]
    mode = sys.argv[2].lower()

    if not os.path.isfile(file_path):
        print(
            json.dumps({"error": f"文件不存在: {file_path}"}, ensure_ascii=False),
            file=sys.stderr,
        )
        sys.exit(1)

    file_size = os.path.getsize(file_path)
    if file_size == 0:
        print(
            json.dumps({"error": "文件大小为 0，请检查文件是否上传完整"}, ensure_ascii=False),
            file=sys.stderr,
        )
        sys.exit(1)
    if file_size < 100 and mode in ("pdf", "pptx", "docx", "doc", "xlsx"):
        print(
            json.dumps(
                {"error": f"文件大小异常（{file_size} 字节），可能不是有效文件"},
                ensure_ascii=False,
            ),
            file=sys.stderr,
        )
        sys.exit(1)

    try:
        if mode == "pptx":
            out = extract_pptx(file_path)
        elif mode == "pdf":
            out = extract_pdf(file_path)
        elif mode == "docx":
            out = extract_docx(file_path)
        elif mode == "doc":
            out = extract_legacy_doc(file_path)
        elif mode == "xlsx":
            out = extract_xlsx(file_path)
        elif mode == "csv":
            out = extract_csv(file_path)
        else:
            raise ValueError(f"不支持的 mode: {mode}，请使用 pptx、pdf、docx、doc、xlsx 或 csv")

        if not out:
            print("警告: 未能从文件中提取到任何文本，返回兜底文本", file=sys.stderr)
            out = extraction_fallback_text(file_path, mode, "未能从文件中提取到任何文本。")

        print(out, end="")

    except Exception as e:
        print(json.dumps({"warning": str(e)}, ensure_ascii=False), file=sys.stderr)
        print(extraction_fallback_text(file_path, mode, str(e)), end="")


if __name__ == "__main__":
    main()
