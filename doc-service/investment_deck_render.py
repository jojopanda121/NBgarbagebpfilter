"""
investment_deck_render.py
─────────────────────────
可变页数投决材料 deck · 确定性渲染器

模型只产 content_schema.json 约束下的内容; 本文件锁定所有版式/颜色/字号.
"""

import json
import sys

from pptx import Presentation
from pptx.util import Inches

from brand_tokens import (
    COLOR, FONT_CN_SERIF, SIZE, HAIRLINE_PT,
    add_rect, set_run, add_text, add_para,
)


SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
LEFT = 0.55
RIGHT = 12.78
CONTENT_W = RIGHT - LEFT


def _clip(value, n=80, fallback="未披露"):
    text = " ".join(str(value or "").split()).strip()
    if not text:
        return fallback
    return text if len(text) <= n else text[: n - 1] + "…"


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR["bg"]
    return slide


def _txt(slide, geom, text, *, size=SIZE["body"], bold=False, color=None,
         align="l", anchor="t", margin=0.05, serif=False):
    tf = add_text(slide, geom, align=align, anchor=anchor, margin=margin)
    p = add_para(tf, align)
    kwargs = {"size": size, "bold": bold, "color": color or COLOR["ink"]}
    if serif:
        kwargs["font_cn"] = FONT_CN_SERIF
    set_run(p.add_run(), text, **kwargs)
    return tf


def _page_header(slide, slide_data, total, company):
    section = _clip(slide_data.get("section_title"), 18, "投决材料")
    title = _clip(slide_data.get("title"), 34, "本页标题")
    add_rect(slide, {"x": 0, "y": 0, "w": SLIDE_W_IN, "h": 0.56}, fill=COLOR["navy"])
    _txt(slide, {"x": LEFT, "y": 0.12, "w": 2.5, "h": 0.28}, section,
         size=9, bold=True, color=COLOR["white"], anchor="m", margin=0)
    _txt(slide, {"x": 3.0, "y": 0.09, "w": 8.4, "h": 0.34}, title,
         size=15, bold=True, color=COLOR["white"], anchor="m", margin=0)
    _txt(slide, {"x": 11.7, "y": 0.12, "w": 1.08, "h": 0.28},
         f"{slide_data.get('page_no', 1)} / {total}", size=8,
         color=COLOR["bg4"], align="r", anchor="m", margin=0)
    add_rect(slide, {"x": LEFT, "y": 0.70, "w": CONTENT_W, "h": 0.045}, fill=COLOR["accent"])
    _txt(slide, {"x": LEFT, "y": 7.12, "w": CONTENT_W, "h": 0.22},
         f"{_clip(company, 24)} · {_clip(slide_data.get('source_note'), 58, '来源: 材料未披露, 待核实')}",
         size=7.5, color=COLOR["mid"], anchor="b", margin=0)


def _insight(slide, text, y=0.86):
    add_rect(slide, {"x": LEFT, "y": y, "w": CONTENT_W, "h": 0.50},
             fill=COLOR["bg3"], line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    _txt(slide, {"x": LEFT + 0.14, "y": y + 0.08, "w": CONTENT_W - 0.28, "h": 0.34},
         _clip(text, 92), size=12, bold=True, color=COLOR["navy"], anchor="m")


def _block_card(slide, geom, block, idx=0, risk=False):
    fill = COLOR["red_bg"] if risk else COLOR["bg2"]
    line = COLOR["red"] if risk else COLOR["border"]
    add_rect(slide, geom, fill=fill, line_color=line, line_w_pt=HAIRLINE_PT)
    label = _clip(block.get("label"), 14)
    value = _clip(block.get("value"), 18, "")
    text = _clip(block.get("text"), 70)
    _txt(slide, {"x": geom["x"] + 0.14, "y": geom["y"] + 0.10, "w": geom["w"] - 0.28, "h": 0.25},
         label, size=10, bold=True, color=COLOR["red"] if risk else COLOR["accent"], margin=0.02)
    if value:
        _txt(slide, {"x": geom["x"] + 0.14, "y": geom["y"] + 0.42, "w": geom["w"] - 0.28, "h": 0.38},
             value, size=17 if len(value) <= 10 else 13, bold=True, color=COLOR["navy"], margin=0.02)
        text_y = geom["y"] + 0.88
        text_h = geom["h"] - 0.96
    else:
        text_y = geom["y"] + 0.42
        text_h = geom["h"] - 0.50
    _txt(slide, {"x": geom["x"] + 0.14, "y": text_y, "w": geom["w"] - 0.28, "h": max(text_h, 0.25)},
         text, size=8.6, color=COLOR["ink"], margin=0.02)


def _grid_blocks(slide, blocks, x, y, w, h, cols=2, risk=False):
    blocks = (blocks or [])[:6]
    rows = max(1, (len(blocks) + cols - 1) // cols)
    gap = 0.14
    card_w = (w - gap * (cols - 1)) / cols
    card_h = (h - gap * (rows - 1)) / rows
    for i, block in enumerate(blocks):
        col = i % cols
        row = i // cols
        _block_card(slide, {
            "x": x + col * (card_w + gap),
            "y": y + row * (card_h + gap),
            "w": card_w,
            "h": card_h,
        }, block, i, risk=risk)


def _draw_table(slide, table, geom, *, risk=False):
    headers = (table or {}).get("headers") or ["事项", "说明"]
    rows = (table or {}).get("rows") or []
    headers = headers[:5]
    rows = rows[:8]
    cols = len(headers)
    row_h = geom["h"] / max(len(rows) + 1, 2)
    col_w = geom["w"] / cols
    add_rect(slide, geom, fill=COLOR["bg2"], line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    for c, head in enumerate(headers):
        add_rect(slide, {"x": geom["x"] + c * col_w, "y": geom["y"], "w": col_w, "h": row_h},
                 fill=COLOR["red"] if risk else COLOR["navy"])
        _txt(slide, {"x": geom["x"] + c * col_w + 0.04, "y": geom["y"] + 0.04, "w": col_w - 0.08, "h": row_h - 0.08},
             _clip(head, 10), size=8.5, bold=True, color=COLOR["white"], anchor="m")
    for r, row in enumerate(rows):
        y = geom["y"] + row_h * (r + 1)
        add_rect(slide, {"x": geom["x"], "y": y, "w": geom["w"], "h": row_h},
                 fill=COLOR["bg3"] if r % 2 else COLOR["bg2"])
        values = list(row)[:cols]
        while len(values) < cols:
            values.append("")
        for c, value in enumerate(values):
            _txt(slide, {"x": geom["x"] + c * col_w + 0.04, "y": y + 0.03, "w": col_w - 0.08, "h": row_h - 0.06},
                 _clip(value, 28), size=7.6, color=COLOR["ink"], anchor="m")


def _draw_bar_chart(slide, chart, geom):
    labels = (chart or {}).get("labels") or []
    values = (chart or {}).get("values") or []
    n = min(len(labels), len(values), 6)
    if n < 2:
        return _grid_blocks(slide, [], geom["x"], geom["y"], geom["w"], geom["h"])
    add_rect(slide, geom, fill=COLOR["bg2"], line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    max_h = geom["h"] - 0.70
    bar_w = (geom["w"] - 0.7) / n
    for i in range(n):
      x = geom["x"] + 0.35 + i * bar_w
      h = max_h * (0.35 + 0.55 * (i + 1) / n)
      add_rect(slide, {"x": x + bar_w * 0.18, "y": geom["y"] + 0.25 + max_h - h, "w": bar_w * 0.48, "h": h},
               fill=COLOR["accent"] if i == n - 1 else COLOR["bg4"])
      _txt(slide, {"x": x, "y": geom["y"] + 0.15 + max_h - h, "w": bar_w * 0.85, "h": 0.18},
           _clip(values[i], 8), size=7.5, bold=True, color=COLOR["navy"], align="c", margin=0)
      _txt(slide, {"x": x, "y": geom["y"] + geom["h"] - 0.34, "w": bar_w * 0.85, "h": 0.20},
           _clip(labels[i], 8), size=7, color=COLOR["mid"], align="c", margin=0)


def _draw_flow(slide, blocks, geom):
    blocks = (blocks or [])[:5]
    n = max(len(blocks), 3)
    gap = 0.16
    box_w = (geom["w"] - gap * (n - 1)) / n
    for i in range(n):
        block = blocks[i] if i < len(blocks) else {"label": f"环节{i + 1}", "value": "", "text": "待补充"}
        x = geom["x"] + i * (box_w + gap)
        _block_card(slide, {"x": x, "y": geom["y"], "w": box_w, "h": geom["h"]}, block, i)
        if i < n - 1:
            _txt(slide, {"x": x + box_w - 0.02, "y": geom["y"] + geom["h"] / 2 - 0.12, "w": gap + 0.04, "h": 0.24},
                 "→", size=14, bold=True, color=COLOR["accent"], align="c", margin=0)


def _render_section_divider(slide, data, total, company):
    add_rect(slide, {"x": 0, "y": 0, "w": SLIDE_W_IN, "h": SLIDE_H_IN}, fill=COLOR["navy"])
    _txt(slide, {"x": 0.85, "y": 1.40, "w": 11.65, "h": 0.35},
         _clip(data.get("section_title"), 22), size=13, color=COLOR["bg4"], margin=0)
    _txt(slide, {"x": 0.85, "y": 1.86, "w": 11.65, "h": 0.90},
         _clip(data.get("title"), 34), size=30, bold=True, color=COLOR["white"], serif=True, margin=0)
    add_rect(slide, {"x": 0.85, "y": 2.90, "w": 5.2, "h": 0.06}, fill=COLOR["accent"])
    _txt(slide, {"x": 0.85, "y": 3.20, "w": 10.8, "h": 0.54},
         _clip(data.get("insight"), 80), size=14, bold=True, color=COLOR["white"], margin=0)
    _grid_blocks(slide, data.get("blocks"), 0.85, 4.15, 11.65, 1.65, cols=3)
    _txt(slide, {"x": 0.85, "y": 6.94, "w": 11.65, "h": 0.22},
         f"{_clip(company, 24)} · {data.get('page_no', 1)} / {total}",
         size=8, color=COLOR["bg4"], margin=0)


def _render_generic(slide, data, total, company):
    _page_header(slide, data, total, company)
    _insight(slide, data.get("insight"))
    tmpl = data.get("template")
    blocks = data.get("blocks") or []

    if tmpl == "exec_summary_4q":
        _grid_blocks(slide, blocks[:4], LEFT, 1.60, CONTENT_W, 4.95, cols=2)
    elif tmpl == "market_size_chart":
        _grid_blocks(slide, blocks[:3], LEFT, 1.60, 4.25, 4.95, cols=1)
        _draw_bar_chart(slide, data.get("chart"), {"x": 5.05, "y": 1.60, "w": 7.73, "h": 4.95})
    elif tmpl == "value_chain_map":
        _draw_flow(slide, blocks, {"x": LEFT, "y": 2.05, "w": CONTENT_W, "h": 2.25})
        _grid_blocks(slide, blocks[:3], LEFT, 4.55, CONTENT_W, 1.70, cols=3)
    elif tmpl == "competition_matrix":
        if data.get("table"):
            _draw_table(slide, data.get("table"), {"x": LEFT, "y": 1.60, "w": CONTENT_W, "h": 4.95})
        else:
            _grid_blocks(slide, blocks[:4], LEFT, 1.60, CONTENT_W, 4.95, cols=2)
    elif tmpl == "timeline":
        _draw_flow(slide, blocks, {"x": LEFT, "y": 2.20, "w": CONTENT_W, "h": 2.10})
        _grid_blocks(slide, blocks[:4], LEFT, 4.70, CONTENT_W, 1.45, cols=4)
    elif tmpl == "financial_table":
        _draw_table(slide, data.get("table"), {"x": LEFT, "y": 1.60, "w": 7.6, "h": 4.95})
        _grid_blocks(slide, blocks[:3], 8.45, 1.60, 4.33, 4.95, cols=1)
    elif tmpl == "valuation_sensitivity":
        _grid_blocks(slide, blocks[:4], LEFT, 1.60, 4.8, 4.95, cols=1)
        _draw_table(slide, data.get("table"), {"x": 5.55, "y": 1.60, "w": 7.23, "h": 4.95})
    elif tmpl == "risk_mitigation":
        if data.get("table"):
            _draw_table(slide, data.get("table"), {"x": LEFT, "y": 1.60, "w": CONTENT_W, "h": 4.95}, risk=True)
        else:
            _grid_blocks(slide, blocks[:6], LEFT, 1.60, CONTENT_W, 4.95, cols=3, risk=True)
    elif tmpl == "next_steps":
        _grid_blocks(slide, blocks[:6], LEFT, 1.60, CONTENT_W, 4.95, cols=3)
    else:
        left_w = 5.55
        _grid_blocks(slide, blocks[:4], LEFT, 1.60, left_w, 4.95, cols=1)
        if data.get("table"):
            _draw_table(slide, data.get("table"), {"x": 6.35, "y": 1.60, "w": 6.43, "h": 4.95})
        elif data.get("chart"):
            _draw_bar_chart(slide, data.get("chart"), {"x": 6.35, "y": 1.60, "w": 6.43, "h": 4.95})
        else:
            _grid_blocks(slide, blocks[4:6] or blocks[:2], 6.35, 1.60, 6.43, 4.95, cols=1)


def render(content, out_path):
    prs = _new_prs()
    company = content.get("company_full_name", "未命名公司")
    slides = content.get("slides") or []
    total = len(slides)
    for slide_data in slides:
        slide = _new_slide(prs)
        if slide_data.get("template") == "section_divider":
            _render_section_divider(slide, slide_data, total, company)
        else:
            _render_generic(slide, slide_data, total, company)
    prs.save(out_path)


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: investment_deck_render.py input.json output.pptx")
        sys.exit(1)
    with open(sys.argv[1], "r", encoding="utf-8") as f:
        data = json.load(f)
    render(data, sys.argv[2])
