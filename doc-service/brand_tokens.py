"""
brand_tokens.py — Python 渲染端的设计令牌单一真源

与 client/src/index.css 的 :root 完全对齐。任何 PPT/PDF render 脚本一律
`from brand_tokens import COLOR, FONT_*, SIZE, set_run, add_rect ...`,
禁止在 render 脚本里硬编码任何颜色 / 字体 / 字号。

修改本文件时, 必须同步更新 client/src/index.css 的 :root 块。
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ═══════════════════════════════════════════════════════════════════
# 颜色 — 与 client/src/index.css :root 同源
# ═══════════════════════════════════════════════════════════════════
COLOR = {
    # 主色族 (权威/品牌)
    "navy":     RGBColor(0x0D, 0x21, 0x45),  # --navy   深海军蓝, 主深底 / 封面 / 标题横幅
    "navy2":    RGBColor(0x16, 0x30, 0x69),  # --navy2  次深底, 渐变 / 副横幅
    "accent":   RGBColor(0x1B, 0x4F, 0xD8),  # --accent 品牌蓝, 强调线 / 内联小标 / KPI 数字
    "accent2":  RGBColor(0x3B, 0x6E, 0xF5),  # --accent2 亮蓝, hover / 副强调
    "gold":     RGBColor(0xA0, 0x70, 0x0A),  # --gold   点睛细线, 用量克制

    # 文本族
    "ink":      RGBColor(0x0F, 0x1C, 0x36),  # --text  主正文
    "mid":      RGBColor(0x4B, 0x5A, 0x72),  # --mid   次级文本 / 表头
    "dim":      RGBColor(0x8E, 0x9B, 0xB0),  # --dim   占位 / 说明

    # 背景族
    "bg":       RGBColor(0xF6, 0xF7, 0xFA),  # --bg    页底浅灰
    "bg2":      RGBColor(0xFF, 0xFF, 0xFF),  # --bg2   卡片白
    "bg3":      RGBColor(0xEE, 0xF1, 0xF7),  # --bg3   信息条带 / 段标签
    "bg4":      RGBColor(0xE5, 0xE9, 0xF4),  # --bg4   更浅信息条

    # 边框族
    "border":   RGBColor(0xD8, 0xDC, 0xE8),  # --border 主边框
    "border2":  RGBColor(0xBF, 0xC5, 0xD6),  # --border2 次级边框 (替代金色虚线)

    # 语义色
    "red":      RGBColor(0xB9, 0x1C, 0x1C),  # --red    风险红
    "red_bg":   RGBColor(0xFE, 0xF2, 0xF2),  # --red-bg 风险浅底
    "green":    RGBColor(0x15, 0x80, 0x3D),  # --green  通过绿
    "green_bg": RGBColor(0xF0, 0xFD, 0xF4),
    "amber":    RGBColor(0xB4, 0x53, 0x09),  # --amber  警告琥珀
    "amber_bg": RGBColor(0xFF, 0xFB, 0xEB),

    "white":    RGBColor(0xFF, 0xFF, 0xFF),
}


# ═══════════════════════════════════════════════════════════════════
# 字体 — 与 client --serif / --sans / --mono 对齐
# PPT 不能嵌入字体, 选系统最大公约数; CN 端用思源宋体, 回退到 PingFang/微软雅黑.
# ═══════════════════════════════════════════════════════════════════
FONT_CN_SERIF = "Noto Serif CJK SC"  # 中文衬线 (标题/封面)
FONT_CN_SANS  = "PingFang SC"         # 中文无衬线 (正文, macOS 优先)
FONT_EN       = "DM Sans"             # 西文/数字 (与 web --sans 对齐)
FONT_MONO     = "JetBrains Mono"      # 等宽 (表格数字, 与 web --mono 对齐)


# ═══════════════════════════════════════════════════════════════════
# 字号 — 与 web 字号阶梯对齐 (PPT pt 与 web px 1:1 近似)
# ═══════════════════════════════════════════════════════════════════
SIZE = {
    "cover_title":  30,
    "cover_tag":    16,
    "cover_meta":   12,
    "page_title":   20,
    "thesis":       14,
    "section":      13,
    "body":         11,
    "table":        10,
    "label_inline": 11,
    "footer":        9,
}


# ═══════════════════════════════════════════════════════════════════
# 装饰常量 — PE/VC 报告倾向冷峻直角, 不要"科技感圆角"
# ═══════════════════════════════════════════════════════════════════
RULE_PT     = 1.25   # 标题下细线粗细
HAIRLINE_PT = 0.5    # 卡片细边粗细


# ═══════════════════════════════════════════════════════════════════
# 共享渲染工具 — 所有 render 脚本必须从这里 import, 不要自己复制实现
# ═══════════════════════════════════════════════════════════════════
def add_rect(slide, geom, fill=None, line_color=None, line_w_pt=None, dash=False):
    """添加一个矩形 (带或不带填充/边框/虚线)"""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(geom["x"]), Inches(geom["y"]),
        Inches(geom["w"]), Inches(geom["h"]),
    )
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        if line_w_pt:
            shp.line.width = Pt(line_w_pt)
        if dash:
            ln = shp.line._get_or_add_ln()
            prstDash = etree.SubElement(ln, qn("a:prstDash"))
            prstDash.set("val", "dash")
    return shp


def set_run(run, text, *, size, bold=False, color=None,
            font_cn=FONT_CN_SANS, font_en=FONT_EN):
    """设置一个 run 的样式 — 中文走 a:ea, 西文走 a:latin."""
    if color is None:
        color = COLOR["ink"]
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", font_en)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_cn)


def add_text(slide, geom, *, align="l", anchor="t", margin=0.05):
    """添加一个空 textbox, 返回 text_frame 让调用方填段落."""
    tb = slide.shapes.add_textbox(
        Inches(geom["x"]), Inches(geom["y"]),
        Inches(geom["w"]), Inches(geom["h"]),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin * 0.6)
    tf.margin_bottom = Inches(margin * 0.6)
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}[anchor]
    tf.paragraphs[0].alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
    return tf


def add_para(tf, align="l", space_after_pt=2):
    """追加一个段落 (首段直接用 paragraphs[0], 其余用 add_paragraph)."""
    if not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
    p.space_after = Pt(space_after_pt)
    return p
