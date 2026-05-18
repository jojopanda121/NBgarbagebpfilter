"""
investment_snapshot_render.py (PE/VC 版)
─────────────────────────────────────────
一页纸投决速览 · 一级市场视角 · 确定性渲染器

用法:
    python investment_snapshot_render.py content.json output.pptx

设计原则:
    · 7 模块版式 / 坐标 / 颜色 / 字号全部锁死.
    · 颜色 / 字体 / 字号一律从 brand_tokens 导入, 与网页 :root 同源.
    · Agent 只产符合 content_schema.json (PE/VC 版) 的 JSON.
    · 同一份 JSON 必须产出结构等价的 PPT (见 brandConsistency.test.js).

7 模块布局:
    Top:        company_full_name + stage_tag chip
                accent rule
                thesis (navy banner)
    Row 1:      [Left 公司概况 summary]    [Right 本轮交易 pe_snapshot 2x2]
    Row 2:      [Left 主要股东 shareholders] [Right 牵引指标 traction 3-row]
    Row 3:      投资亮点 highlights × 4
    Row 4:      投资风险 risks × 2 (含 mitigant 缓解)
"""

import json
import sys

from pptx import Presentation
from pptx.util import Inches

from brand_tokens import (
    COLOR, FONT_CN_SERIF, FONT_CN_SANS, FONT_EN, SIZE,
    RULE_PT, HAIRLINE_PT,
    add_rect, set_run, add_text, add_para,
)


# 画布: A4 横版
SLIDE_W = Inches(11.69)
SLIDE_H = Inches(8.27)

# 7 模块坐标 (英寸). 重排后行高比原版略调以容纳新模块.
LAYOUT = {
    # 顶部
    "title":         {"x": 0.82, "y": 0.42, "w": 8.50, "h": 0.40},
    "stage_chip":    {"x": 9.50, "y": 0.46, "w": 1.40, "h": 0.32},
    "accent_rule":   {"x": 0.82, "y": 0.86, "w": 10.05, "h": 0.05},
    "thesis_banner": {"x": 0.82, "y": 1.10, "w": 10.05, "h": 0.36},

    # Row 1: 公司概况 | 本轮交易 (上半 1.20")
    "card_overview":   {"x": 0.82, "y": 1.74, "w": 4.92, "h": 1.20},
    "label_overview":  {"x": 2.32, "y": 1.52, "w": 1.92, "h": 0.30},
    "body_overview":   {"x": 0.96, "y": 1.86, "w": 4.63, "h": 1.06},

    "card_pe":         {"x": 5.96, "y": 1.74, "w": 4.92, "h": 1.20},
    "label_pe":        {"x": 7.45, "y": 1.52, "w": 1.92, "h": 0.30},
    "body_pe":         {"x": 6.05, "y": 1.86, "w": 4.71, "h": 1.06},

    # Row 2: 主要股东 | 牵引指标 (下半 1.30")
    "card_share":      {"x": 0.82, "y": 3.04, "w": 4.92, "h": 1.30},
    "label_share":     {"x": 2.32, "y": 2.82, "w": 1.92, "h": 0.30},
    "body_share":      {"x": 0.96, "y": 3.16, "w": 4.63, "h": 1.16},

    "card_trac":       {"x": 5.96, "y": 3.04, "w": 4.92, "h": 1.30},
    "label_trac":      {"x": 7.45, "y": 2.82, "w": 1.92, "h": 0.30},
    "body_trac":       {"x": 6.05, "y": 3.16, "w": 4.71, "h": 1.16},

    # Row 3: 投资亮点 (满宽)
    "card_high":     {"x": 0.82, "y": 4.62, "w": 10.05, "h": 2.06},
    "label_high":    {"x": 4.86, "y": 4.40, "w": 1.98, "h": 0.30},
    "body_high":     {"x": 0.96, "y": 4.76, "w": 9.77, "h": 1.86},

    # Row 4: 投资风险 (满宽, 含 mitigant)
    "card_risk":     {"x": 0.82, "y": 6.90, "w": 10.05, "h": 1.10},
    "label_risk":    {"x": 0.84, "y": 6.92, "w": 1.52, "h": 1.06},
    "block_risk":    {"x": 2.36, "y": 6.92, "w": 8.49, "h": 1.06},
    "body_risk":     {"x": 2.44, "y": 6.96, "w": 8.35, "h": 0.98},
}


def render(content: dict, out_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ─── 标题 + 阶段 chip ───
    tf = add_text(slide, LAYOUT["title"], align="l", margin=0)
    p = add_para(tf, "l")
    set_run(p.add_run(), f"投资要点速览——{content['company_full_name']}",
            size=SIZE["page_title"], bold=True, color=COLOR["navy"],
            font_cn=FONT_CN_SERIF)
    # stage_tag chip (右上角胶囊, navy 描边 + accent 文字)
    add_rect(slide, LAYOUT["stage_chip"],
             fill=COLOR["bg3"], line_color=COLOR["accent"], line_w_pt=HAIRLINE_PT)
    tf = add_text(slide, LAYOUT["stage_chip"], align="c", anchor="m", margin=0)
    p = add_para(tf, "c")
    set_run(p.add_run(), content["stage_tag"],
            size=SIZE["label_inline"], bold=True, color=COLOR["accent"])

    # ─── 品牌蓝细线 ───
    add_rect(slide, LAYOUT["accent_rule"], fill=COLOR["accent"])

    # ─── 论点 navy 横幅 ───
    add_rect(slide, LAYOUT["thesis_banner"], fill=COLOR["navy"])
    tf = add_text(slide, LAYOUT["thesis_banner"], align="c", anchor="m", margin=0.1)
    p = add_para(tf, "c")
    set_run(p.add_run(), content["thesis"],
            size=SIZE["thesis"], bold=True, color=COLOR["white"])

    # ─── Row 1 · 公司概况 (左上) ───
    add_rect(slide, LAYOUT["card_overview"], fill=COLOR["bg2"],
             line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    _draw_section_label(slide, LAYOUT["label_overview"], "公司概况")
    tf = add_text(slide, LAYOUT["body_overview"], margin=0.08)
    p = add_para(tf, "l", space_after_pt=2)
    set_run(p.add_run(), content["company_overview"]["summary"],
            size=SIZE["body"], color=COLOR["ink"])

    # ─── Row 1 · 本轮交易快照 (右上) — 2x2 KPI 网格 ───
    add_rect(slide, LAYOUT["card_pe"], fill=COLOR["bg2"],
             line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    _draw_section_label(slide, LAYOUT["label_pe"], "本轮交易")
    _draw_pe_snapshot(slide, content["pe_snapshot"])

    # ─── Row 2 · 主要股东 (左下) ───
    add_rect(slide, LAYOUT["card_share"], fill=COLOR["bg2"],
             line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    _draw_section_label(slide, LAYOUT["label_share"], "主要股东")
    _draw_shareholders(slide, content["cap_and_traction"]["shareholders"])

    # ─── Row 2 · 牵引指标 (右下) ───
    add_rect(slide, LAYOUT["card_trac"], fill=COLOR["bg2"],
             line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    _draw_section_label(slide, LAYOUT["label_trac"], "牵引指标")
    _draw_traction(slide, content["cap_and_traction"]["traction"])

    # ─── Row 3 · 投资亮点 ───
    add_rect(slide, LAYOUT["card_high"], fill=COLOR["bg2"],
             line_color=COLOR["border"], line_w_pt=HAIRLINE_PT)
    _draw_section_label(slide, LAYOUT["label_high"], "投资亮点")
    tf = add_text(slide, LAYOUT["body_high"], margin=0.08)
    for item in content["highlights"]:
        p = add_para(tf, "l", space_after_pt=3)
        set_run(p.add_run(), f"{item['label']}：",
                size=SIZE["label_inline"], bold=True, color=COLOR["accent"])
        set_run(p.add_run(), item["desc"],
                size=SIZE["body"], color=COLOR["ink"])

    # ─── Row 4 · 投资风险 (含 mitigant) ───
    add_rect(slide, LAYOUT["card_risk"], fill=COLOR["bg2"],
             line_color=COLOR["red"], line_w_pt=HAIRLINE_PT)
    add_rect(slide, LAYOUT["label_risk"], fill=COLOR["red"])
    tf = add_text(slide, LAYOUT["label_risk"], align="c", anchor="m", margin=0)
    p = add_para(tf, "c")
    set_run(p.add_run(), "投资风险",
            size=SIZE["section"], bold=True, color=COLOR["white"])
    add_rect(slide, LAYOUT["block_risk"], fill=COLOR["red_bg"])
    tf = add_text(slide, LAYOUT["body_risk"], margin=0.08)
    # Deal-breakers 先行 (如果有). 与 risks 视觉区分: "致命伤 ⚠️" 加深红粗体 + 证伪条件.
    deal_breakers = content.get("deal_breakers") or []
    for db in deal_breakers:
        p = add_para(tf, "l", space_after_pt=2)
        set_run(p.add_run(), f"致命伤 · {db.get('title', '')}：",
                size=SIZE["label_inline"], bold=True, color=COLOR["red"])
        set_run(p.add_run(), db.get("logic", ""),
                size=SIZE["body"], color=COLOR["ink"])
        p = add_para(tf, "l", space_after_pt=2)
        set_run(p.add_run(), "证伪条件：",
                size=SIZE["label_inline"], bold=True, color=COLOR["mid"])
        set_run(p.add_run(), db.get("falsification_test", ""),
                size=SIZE["body"], color=COLOR["mid"])
    for item in content["risks"]:
        p = add_para(tf, "l", space_after_pt=2)
        set_run(p.add_run(), f"{item['label']}：",
                size=SIZE["label_inline"], bold=True, color=COLOR["red"])
        set_run(p.add_run(), item["desc"],
                size=SIZE["body"], color=COLOR["ink"])
        # 缓解措施: 用 mid 色, 与正文区分; 前缀 "缓解: "
        p = add_para(tf, "l", space_after_pt=2)
        set_run(p.add_run(), "缓解：",
                size=SIZE["label_inline"], bold=True, color=COLOR["mid"])
        set_run(p.add_run(), item["mitigant"],
                size=SIZE["body"], color=COLOR["mid"])

    prs.save(out_path)


def _draw_section_label(slide, geom, text):
    """段标签: 浅蓝灰胶囊 + 品牌蓝文字, 与网页 status pill 一致."""
    add_rect(slide, geom, fill=COLOR["bg3"])
    tf = add_text(slide, geom, align="c", anchor="m", margin=0)
    p = add_para(tf, "c")
    set_run(p.add_run(), text,
            size=SIZE["section"], bold=True, color=COLOR["accent"])


def _draw_pe_snapshot(slide, pe):
    """本轮交易 2x2 KPI 网格: 阶段 / 本轮规模 / Pre 估值 / 关键条款"""
    o = LAYOUT["body_pe"]
    # 2 列 × 2 行, 上一行: stage / round_size; 下一行: pre_valuation / lead_terms
    col_w = o["w"] / 2
    row_h = o["h"] / 2
    cells = [
        ("阶段",       pe["stage"]),
        ("本轮规模",   pe["round_size"]),
        ("Pre 估值",   pe["pre_valuation"]),
        ("关键条款",   pe["lead_terms"]),
    ]
    for i, (label, value) in enumerate(cells):
        col, row = i % 2, i // 2
        x = o["x"] + col * col_w
        y = o["y"] + row * row_h
        # 标签 (上半)
        tf = add_text(slide, {"x": x, "y": y, "w": col_w, "h": row_h * 0.40},
                      align="l", anchor="t", margin=0.04)
        set_run(add_para(tf, "l").add_run(), label,
                size=SIZE["table"], bold=True, color=COLOR["mid"])
        # 值 (下半)
        tf = add_text(slide,
                      {"x": x, "y": y + row_h * 0.40,
                       "w": col_w, "h": row_h * 0.60},
                      align="l", anchor="t", margin=0.04)
        set_run(add_para(tf, "l").add_run(), value,
                size=SIZE["body"], bold=True, color=COLOR["navy"])


def _draw_shareholders(slide, shareholders):
    """主要股东表: 2 列 (name / pct), 行高自适应, 顶下分隔线品牌蓝."""
    o = LAYOUT["body_share"]
    n = len(shareholders)
    name_w = o["w"] * 0.70
    pct_w = o["w"] * 0.30
    row_h = o["h"] / max(n, 1)
    for i, sh in enumerate(shareholders):
        y = o["y"] + i * row_h
        # name (左对齐)
        tf = add_text(slide, {"x": o["x"], "y": y, "w": name_w, "h": row_h},
                      align="l", anchor="m", margin=0.04)
        set_run(add_para(tf, "l").add_run(), sh["name"],
                size=SIZE["table"], color=COLOR["ink"])
        # pct (右对齐, accent 色)
        tf = add_text(slide,
                      {"x": o["x"] + name_w, "y": y, "w": pct_w, "h": row_h},
                      align="r", anchor="m", margin=0.04)
        set_run(add_para(tf, "r").add_run(), sh["pct"],
                size=SIZE["table"], bold=True, color=COLOR["accent"])
        # 行间细分隔线 (除最后一行)
        if i < n - 1:
            add_rect(slide,
                     {"x": o["x"], "y": y + row_h - 0.01,
                      "w": o["w"], "h": 0.008},
                     fill=COLOR["border"])


def _draw_traction(slide, traction):
    """牵引指标 3 行: metric / value / note. 与股东表同款行高."""
    o = LAYOUT["body_trac"]
    metric_w = o["w"] * 0.22
    value_w = o["w"] * 0.28
    note_w = o["w"] * 0.50
    row_h = o["h"] / max(len(traction), 1)
    for i, t in enumerate(traction):
        y = o["y"] + i * row_h
        # metric (mid 色)
        tf = add_text(slide, {"x": o["x"], "y": y, "w": metric_w, "h": row_h},
                      align="l", anchor="m", margin=0.04)
        set_run(add_para(tf, "l").add_run(), t["metric"],
                size=SIZE["table"], bold=True, color=COLOR["mid"])
        # value (navy bold, 主数据)
        tf = add_text(slide,
                      {"x": o["x"] + metric_w, "y": y,
                       "w": value_w, "h": row_h},
                      align="l", anchor="m", margin=0.04)
        set_run(add_para(tf, "l").add_run(), t["value"],
                size=SIZE["body"], bold=True, color=COLOR["navy"])
        # note (ink, 对标说明)
        tf = add_text(slide,
                      {"x": o["x"] + metric_w + value_w, "y": y,
                       "w": note_w, "h": row_h},
                      align="l", anchor="m", margin=0.04)
        set_run(add_para(tf, "l").add_run(), t["note"],
                size=SIZE["table"], color=COLOR["ink"])
        if i < len(traction) - 1:
            add_rect(slide,
                     {"x": o["x"], "y": y + row_h - 0.01,
                      "w": o["w"], "h": 0.008},
                     fill=COLOR["border"])


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python investment_snapshot_render.py <content.json> <output.pptx>")
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        content = json.load(f)
    render(content, sys.argv[2])
    print(f"✓ rendered → {sys.argv[2]}")
