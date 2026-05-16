"""
doc-service/main.py — 文档提取微服务 (FastAPI)

独立微服务，替代 Node.js 的 child_process 调用方式。
Node.js 主后端通过 HTTP 将文件传给 Python 服务进行解析。

启动方式：
  uvicorn main:app --host 0.0.0.0 --port 8001

Docker 启动：
  docker build -t doc-service .
  docker run -p 8001:8001 doc-service
"""

import io
import csv
import os
import tempfile
from pathlib import Path
from typing import List, Optional

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel

app = FastAPI(title="GarbageBPFilter Doc Service", version="1.0.0")

# 在应用启动时初始化 OCR 引擎（避免每次请求重新加载模型）
_ocr_engine = None


def _get_ocr():
    global _ocr_engine
    if _ocr_engine is None:
        from rapidocr_onnxruntime import RapidOCR
        _ocr_engine = RapidOCR()
    return _ocr_engine


def extract_pdf_text(file_path: str) -> str:
    """从 PDF 提取文本，先尝试直接提取，文本过少时降级为 OCR"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    pages_text = []

    for page in doc:
        text = page.get_text("text")
        pages_text.append(text.strip())

    doc.close()
    full_text = "\n".join(pages_text)

    # 如果每页平均文本不足 50 字符，尝试 OCR
    avg_chars = len(full_text) / max(len(pages_text), 1)
    if avg_chars < 50:
        full_text = extract_pdf_ocr(file_path)

    return full_text


def extract_pdf_ocr(file_path: str) -> str:
    """OCR 提取 PDF 文本"""
    import fitz

    ocr = _get_ocr()
    doc = fitz.open(file_path)
    pages_text = []

    for page in doc:
        pix = page.get_pixmap(dpi=200)
        img_bytes = pix.tobytes("png")

        result, _ = ocr(img_bytes)
        if result:
            page_text = "\n".join([line[1] for line in result])
            pages_text.append(page_text)

    doc.close()
    return "\n".join(pages_text)


def extract_pptx_text(file_path: str) -> str:
    """从 PPTX 提取文本"""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_text = []

    for slide in prs.slides:
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            slides_text.append("\n".join(slide_parts))

    return "\n\n".join(slides_text)


def extract_docx_text(file_path: str) -> str:
    """从 DOCX 提取段落和表格文本"""
    from docx import Document

    doc = Document(file_path)
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    for table_idx, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[表格 {table_idx}]\n" + "\n".join(rows))

    return "\n\n".join(parts)


def extract_xlsx_text(file_path: str) -> str:
    """从 XLSX 提取工作表预览文本，保留表头和前 80 行"""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True, read_only=True)
    parts = []

    for ws in wb.worksheets:
        lines = [f"# Sheet: {ws.title}"]
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if not any(v.strip() for v in values):
                continue
            lines.append(" | ".join(values[:20]))
            row_count += 1
            if row_count >= 80:
                lines.append("...（已截断，仅展示前 80 行非空数据）")
                break
        if row_count > 0:
            parts.append("\n".join(lines))

    wb.close()
    return "\n\n".join(parts)


def extract_csv_text(file_path: str) -> str:
    """从 CSV 提取前 120 行文本"""
    encodings = ("utf-8-sig", "utf-8", "gb18030")
    last_error = None
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc, newline="") as f:
                reader = csv.reader(f)
                lines = []
                for idx, row in enumerate(reader):
                    if idx >= 120:
                        lines.append("...（已截断，仅展示前 120 行）")
                        break
                    lines.append(" | ".join(row[:30]))
                return "\n".join(lines)
        except UnicodeDecodeError as e:
            last_error = e
    raise RuntimeError(f"CSV 编码识别失败: {last_error}")


@app.post("/extract")
async def extract_text(
    file: UploadFile = File(...),
    mode: str = Form("pdf"),
):
    """
    接收上传的文档文件，提取文本内容。

    Args:
        file: 上传的 PDF 或 PPTX 文件
        mode: 文件类型 ("pdf"、"pptx"、"docx"、"xlsx" 或 "csv")

    Returns:
        { "text": "提取的文本内容", "chars": 字符数 }
    """
    if mode not in ("pdf", "pptx", "docx", "xlsx", "csv"):
        raise HTTPException(status_code=400, detail="不支持的文件类型，仅支持 pdf/pptx/docx/xlsx/csv")

    # 保存到临时文件
    suffix = f".{mode}"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        content = await file.read()
        tmp.write(content)
        tmp_path = tmp.name

    try:
        if mode == "pptx":
            text = extract_pptx_text(tmp_path)
        elif mode == "docx":
            text = extract_docx_text(tmp_path)
        elif mode == "xlsx":
            text = extract_xlsx_text(tmp_path)
        elif mode == "csv":
            text = extract_csv_text(tmp_path)
        else:
            text = extract_pdf_text(tmp_path)

        if not text or len(text.strip()) < 10:
            raise HTTPException(
                status_code=422,
                detail="文档提取的文本过少，请检查文件是否包含有效内容",
            )

        return JSONResponse({
            "text": text,
            "chars": len(text),
        })

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文档解析失败: {str(e)}")
    finally:
        # 清理临时文件
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


# ── PPT 生成 ────────────────────────────────────────────────

class SlideSpec(BaseModel):
    title: str
    bullets: Optional[List[str]] = None
    notes: Optional[str] = None


class GeneratePptxPayload(BaseModel):
    title: str
    subtitle: Optional[str] = None
    slides: List[SlideSpec]


def _build_pptx(payload: GeneratePptxPayload) -> bytes:
    """根据 slides JSON 渲染 .pptx 字节流（深色简洁模板）"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG = RGBColor(0x0F, 0x17, 0x2A)
    ACCENT = RGBColor(0x60, 0xA5, 0xFA)
    TEXT = RGBColor(0xE2, 0xE8, 0xF0)
    SUB = RGBColor(0x94, 0xA3, 0xB8)

    def paint_bg(slide):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg_shape.line.fill.background()
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG
        # 把矩形放到最底层
        spTree = bg_shape._element.getparent()
        spTree.remove(bg_shape._element)
        spTree.insert(2, bg_shape._element)

    blank_layout = prs.slide_layouts[6]

    skip_cover = len(payload.slides) == 1 and (payload.subtitle or "").lower() == "one page"
    if not skip_cover:
        # 封面页
        cover = prs.slides.add_slide(blank_layout)
        paint_bg(cover)
        title_box = cover.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = payload.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TEXT

        if payload.subtitle:
            sub_box = cover.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.8))
            sp = sub_box.text_frame.paragraphs[0]
            sp.text = payload.subtitle
            sp.font.size = Pt(20)
            sp.font.color.rgb = SUB

    # 内容页
    for idx, slide_spec in enumerate(payload.slides, 1):
        s = prs.slides.add_slide(blank_layout)
        paint_bg(s)

        # 顶部装饰条
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(0.4), Inches(0.05))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT

        # 标题
        title_box = s.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(1.0))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = slide_spec.title or f"第 {idx} 页"
        tp.font.size = Pt(30)
        tp.font.bold = True
        tp.font.color.rgb = TEXT

        # 项目符号
        bullets = slide_spec.bullets or []
        if bullets:
            body = s.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(5.0))
            tf = body.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = f"•  {bullet}"
                para.font.size = Pt(18)
                para.font.color.rgb = TEXT
                para.space_after = Pt(8)

        # 页脚页码
        pn = s.shapes.add_textbox(Inches(12.4), Inches(7.0), Inches(0.7), Inches(0.4))
        pp = pn.text_frame.paragraphs[0]
        pp.text = f"{idx}/{len(payload.slides)}"
        pp.font.size = Pt(10)
        pp.font.color.rgb = SUB

        if slide_spec.notes:
            s.notes_slide.notes_text_frame.text = slide_spec.notes

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


@app.post("/generate/pptx")
async def generate_pptx(payload: GeneratePptxPayload):
    if not payload.slides:
        raise HTTPException(status_code=400, detail="slides 不能为空")
    try:
        data = _build_pptx(payload)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPT 生成失败: {str(e)}")
    headers = {
        "Content-Disposition": "attachment; filename=presentation.pptx",
    }
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers,
    )


# ── Word / Excel 生成 ───────────────────────────────────────

class DocxSection(BaseModel):
    heading: str
    paragraphs: Optional[List[str]] = None
    bullets: Optional[List[str]] = None


class GenerateDocxPayload(BaseModel):
    title: str
    subtitle: Optional[str] = None
    sections: List[DocxSection]


def _build_docx(payload: GenerateDocxPayload) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    styles = doc.styles
    styles["Normal"].font.name = "Microsoft YaHei"
    styles["Normal"].font.size = Pt(10.5)

    title = doc.add_heading(payload.title, level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if payload.subtitle:
        sub = doc.add_paragraph(payload.subtitle)
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if sub.runs:
            sub.runs[0].font.color.rgb = RGBColor(0x5B, 0x67, 0x7A)

    for section in payload.sections:
        doc.add_heading(section.heading or "未命名章节", level=1)
        for para in section.paragraphs or []:
            if para:
                doc.add_paragraph(para)
        for bullet in section.bullets or []:
            if bullet:
                doc.add_paragraph(bullet, style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


@app.post("/generate/docx")
async def generate_docx(payload: GenerateDocxPayload):
    if not payload.sections:
        raise HTTPException(status_code=400, detail="sections 不能为空")
    try:
        data = _build_docx(payload)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Word 生成失败: {str(e)}")
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": "attachment; filename=document.docx"},
    )


class SheetSpec(BaseModel):
    name: str
    headers: Optional[List[str]] = None
    rows: Optional[List[List[str]]] = None


class GenerateXlsxPayload(BaseModel):
    title: str
    sheets: List[SheetSpec]


def _build_xlsx(payload: GenerateXlsxPayload) -> bytes:
    import xlsxwriter

    buf = io.BytesIO()
    workbook = xlsxwriter.Workbook(buf, {"in_memory": True})
    title_fmt = workbook.add_format({"bold": True, "font_size": 14, "font_color": "#0F1C36"})
    header_fmt = workbook.add_format({"bold": True, "bg_color": "#E5E9F4", "font_color": "#0F1C36", "border": 1})
    cell_fmt = workbook.add_format({"text_wrap": True, "valign": "top", "border": 1})

    for idx, sheet in enumerate(payload.sheets):
        safe_name = (sheet.name or f"Sheet{idx + 1}")[:31].replace("/", "_").replace("\\", "_")
        ws = workbook.add_worksheet(safe_name)
        ws.write(0, 0, payload.title, title_fmt)

        headers = sheet.headers or []
        rows = sheet.rows or []
        start_row = 2
        for col, header in enumerate(headers):
            ws.write(start_row, col, header, header_fmt)
            ws.set_column(col, col, 18)
        for r_idx, row in enumerate(rows, start_row + 1):
            for c_idx, value in enumerate(row):
                ws.write(r_idx, c_idx, value, cell_fmt)
        if headers:
            ws.freeze_panes(start_row + 1, 0)
            ws.autofilter(start_row, 0, start_row + max(len(rows), 1), max(len(headers) - 1, 0))

    workbook.close()
    buf.seek(0)
    return buf.getvalue()


@app.post("/generate/xlsx")
async def generate_xlsx(payload: GenerateXlsxPayload):
    if not payload.sheets:
        raise HTTPException(status_code=400, detail="sheets 不能为空")
    try:
        data = _build_xlsx(payload)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Excel 生成失败: {str(e)}")
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=workbook.xlsx"},
    )


# ── 一页投资亮点 PPT (投资要点速览) ─────────────────────────

class OnePagerProduct(BaseModel):
    name: str
    desc: str


class OnePagerOverview(BaseModel):
    summary: str
    products: List[OnePagerProduct]


class OnePagerKpi(BaseModel):
    label: str
    value: str


class OnePagerDriver(BaseModel):
    type: str
    text: str


class OnePagerMarket(BaseModel):
    kpis: List[OnePagerKpi]
    drivers: List[OnePagerDriver]
    competition: str


class OnePagerItem(BaseModel):
    title: str
    desc: str


class OnePagerFooter(BaseModel):
    founded: str
    team_size: str
    funding_total: str
    ai_grade: str


class OnePagerPayload(BaseModel):
    company_name: str
    headline: str
    company_overview: OnePagerOverview
    market_opportunity: OnePagerMarket
    highlights: List[OnePagerItem]
    risks: List[OnePagerItem]
    footer: OnePagerFooter


def _set_cn_font(run, size_pt: int, bold: bool = False, color=None,
                 family: str = None):
    """设置中文友好字体 (同时设置 latin / eastAsia typeface).

    family 默认走 brand_tokens.FONT_CN_SANS, 与网页 :root 同源.
    """
    from pptx.util import Pt
    from lxml import etree
    from brand_tokens import FONT_CN_SANS, FONT_EN
    if family is None:
        family = FONT_CN_SANS
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = FONT_EN  # latin 走 DM Sans
    if color is not None:
        run.font.color.rgb = color
    # 强制 east-asian / latin / cs typeface, 保证 Mac/Win/WPS 渲染一致
    rPr = run._r.get_or_add_rPr()
    main_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag, typeface in (("ea", family), ("cs", family), ("latin", FONT_EN)):
        existing = rPr.find(f"{main_ns}{tag}")
        if existing is not None:
            rPr.remove(existing)
        el = etree.SubElement(rPr, f"{main_ns}{tag}")
        el.set("typeface", typeface)


def _build_onepager(payload: OnePagerPayload) -> bytes:
    """渲染一页"投资要点速览"PPT.

    版式锁定: navy 横幅 + 品牌蓝细线 + 双栏 (公司概况 / 市场机会)
    + 投资亮点 4 条 + 投资风险 2 条 + 页脚小条.
    颜色 / 字体一律走 brand_tokens, 与网页 :root 同源.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE

    from brand_tokens import COLOR

    # 语义映射 (统一品牌色, 不再使用砖红/米黄)
    BG = COLOR["bg"]            # 页底浅灰
    BRAND = COLOR["accent"]     # 品牌蓝, 标题色 / 分隔线
    BANNER = COLOR["navy"]      # 主横幅 (深海军蓝)
    LABEL_FG = COLOR["accent"]  # 板块标签文字
    BLACK = COLOR["ink"]
    GRAY = COLOR["mid"]
    LIGHT = COLOR["bg3"]        # 板块浅蓝灰底
    GRAY_BG = COLOR["red_bg"]   # 风险板块走语义红浅底
    RISK_FG = COLOR["red"]      # 风险板块文字

    def clip(value, max_chars, fallback="暂无"):
        text = " ".join(str(value or "").split()).strip()
        if not text:
            return fallback
        return text if len(text) <= max_chars else text[: max_chars - 1] + "…"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    def paint_bg():
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, prs.slide_height)
        rect.line.fill.background()
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG
        sp = rect._element.getparent()
        sp.remove(rect._element)
        sp.insert(2, rect._element)

    def add_text(left, top, width, height, text, *, size=12, bold=False,
                 color=BLACK, anchor_top=True, line_spacing=None,
                 fill_color=None, line_color=None):
        box = s.shapes.add_textbox(left, top, width, height)
        if fill_color is not None:
            box.fill.solid()
            box.fill.fore_color.rgb = fill_color
        else:
            box.fill.background()
        if line_color is not None:
            box.line.color.rgb = line_color
            box.line.width = Pt(0.75)
        else:
            box.line.fill.background()
        tf = box.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.TOP if anchor_top else MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text or ""
        _set_cn_font(run, size, bold=bold, color=color)
        return box, tf, p

    def add_filled_rect(left, top, width, height, color):
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.line.fill.background()
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        return rect

    def add_section_label(left, top, width, text):
        """板块小标签 (浅蓝灰胶囊 + 品牌蓝文字, 与网页 status pill 同款)."""
        add_filled_rect(left, top, width, Inches(0.34), LIGHT)
        add_text(left, top, width, Inches(0.34), text,
                 size=14, bold=True, color=LABEL_FG, anchor_top=False)

    paint_bg()

    # ── 顶部标题 + 品牌蓝细线 ───────────────────────────────
    company_name = clip(payload.company_name, 28)
    title_text = f"投资要点速览——{company_name}"
    title_size = 20 if len(company_name) > 22 else 24
    add_text(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.6),
             title_text, size=title_size, bold=True, color=COLOR["navy"])
    add_filled_rect(Inches(0.55), Inches(0.95), Inches(12.2), Inches(0.025), BRAND)

    # ── navy 横幅 (论点带) ──────────────────────────────────
    add_filled_rect(Inches(0.55), Inches(1.10), Inches(12.2), Inches(0.55), BANNER)
    headline_box, _, _ = add_text(
        Inches(0.55), Inches(1.10), Inches(12.2), Inches(0.55),
        clip(payload.headline, 36), size=15, bold=True, color=COLOR["white"],
        anchor_top=False
    )
    headline_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # ── 左：公司概况 ────────────────────────────────────────
    LEFT_X = Inches(0.55)
    LEFT_W = Inches(5.9)
    BLOCK_TOP = Inches(1.78)
    add_section_label(LEFT_X, BLOCK_TOP, LEFT_W, "公司概况")

    # 公司概况框 + 内容（虚线感用浅边框模拟）
    overview_top = Inches(2.14)
    overview_h = Inches(2.25)
    add_filled_rect(LEFT_X, overview_top, LEFT_W, overview_h, BG)

    # 摘要段
    summary_box, summary_tf, _ = add_text(
        LEFT_X, overview_top, LEFT_W, Inches(0.92),
        clip(payload.company_overview.summary, 88), size=9, color=BLACK,
        line_spacing=1.25
    )
    summary_box.line.color.rgb = COLOR["border"]
    summary_box.line.width = Pt(0.5)

    # 产品/业务条目
    p_top = overview_top + Inches(0.98)
    p_h = (overview_h - Inches(0.98)) // max(len(payload.company_overview.products), 1)
    for i, prod in enumerate(payload.company_overview.products[:3]):
        item_top = p_top + p_h * i
        item_box = s.shapes.add_textbox(LEFT_X, item_top, LEFT_W, p_h)
        item_box.line.fill.background()
        item_box.fill.background()
        tf = item_box.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        para = tf.paragraphs[0]
        para.line_spacing = 1.2
        r1 = para.add_run()
        r1.text = f"{clip(prod.name, 10)}："
        _set_cn_font(r1, 9, bold=True, color=LABEL_FG)
        r2 = para.add_run()
        r2.text = clip(prod.desc, 28)
        _set_cn_font(r2, 9, color=BLACK)

    # ── 右：市场机会与行业速览 ─────────────────────────────
    RIGHT_X = Inches(6.65)
    RIGHT_W = Inches(6.10)
    add_section_label(RIGHT_X, BLOCK_TOP, RIGHT_W, "市场机会与行业速览")

    market_top = Inches(2.14)
    market_h = Inches(2.25)
    add_filled_rect(RIGHT_X, market_top, RIGHT_W, market_h, BG)

    # KPI 行（4 列）
    kpi_h = Inches(0.7)
    kpi_w = RIGHT_W // 4
    for i, kpi in enumerate(payload.market_opportunity.kpis[:4]):
        x = RIGHT_X + kpi_w * i
        # 标签
        label_box, _, lp = add_text(
            x, market_top + Inches(0.05), kpi_w, Inches(0.28),
            clip(kpi.label, 8), size=8, bold=True, color=GRAY, anchor_top=False
        )
        lp.alignment = PP_ALIGN.CENTER
        # 值
        val_box, _, vp = add_text(
            x, market_top + Inches(0.30), kpi_w, Inches(0.40),
            clip(kpi.value, 12), size=8, bold=True, color=LABEL_FG, anchor_top=False
        )
        vp.alignment = PP_ALIGN.CENTER

    # 驱动力 3 条
    drv_top = market_top + Inches(0.78)
    drv_h = Inches(0.38)
    for i, drv in enumerate(payload.market_opportunity.drivers[:3]):
        y = drv_top + drv_h * i
        box = s.shapes.add_textbox(RIGHT_X + Inches(0.1), y,
                                    RIGHT_W - Inches(0.2), drv_h)
        box.line.fill.background()
        box.fill.background()
        tf = box.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.line_spacing = 1.15
        r1 = p.add_run()
        r1.text = f"〔{clip(drv.type, 6)}〕 "
        _set_cn_font(r1, 9, bold=True, color=LABEL_FG)
        r2 = p.add_run()
        r2.text = clip(drv.text, 34)
        _set_cn_font(r2, 9, color=BLACK)

    # 竞争格局
    comp_top = drv_top + drv_h * 3 + Inches(0.05)
    comp_box = s.shapes.add_textbox(RIGHT_X + Inches(0.1), comp_top,
                                     RIGHT_W - Inches(0.2), Inches(0.55))
    comp_box.line.fill.background()
    comp_box.fill.background()
    ctf = comp_box.text_frame
    ctf.auto_size = MSO_AUTO_SIZE.NONE
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.05)
    cp = ctf.paragraphs[0]
    cp.line_spacing = 1.15
    cr1 = cp.add_run()
    cr1.text = "〔竞争格局〕 "
    _set_cn_font(cr1, 9, bold=True, color=LABEL_FG)
    cr2 = cp.add_run()
    cr2.text = clip(payload.market_opportunity.competition, 50)
    _set_cn_font(cr2, 9, color=BLACK)

    # ── 投资亮点（4 条，2x2 网格） ─────────────────────────
    HL_TOP = Inches(4.55)
    HL_H = Inches(1.45)
    add_section_label(LEFT_X, HL_TOP, Inches(12.2), "投资亮点")
    cells_top = HL_TOP + Inches(0.40)
    cell_w = (Inches(12.2) - Inches(0.30)) // 2
    cell_h = (HL_H - Inches(0.40)) // 2
    for i, hl in enumerate(payload.highlights[:4]):
        col = i % 2
        row = i // 2
        x = LEFT_X + (cell_w + Inches(0.30)) * col
        y = cells_top + cell_h * row
        # 标题
        add_text(x, y, cell_w, Inches(0.35),
                 f"· {clip(hl.title, 10)}", size=10, bold=True, color=LABEL_FG,
                 anchor_top=False)
        # 描述
        add_text(x + Inches(0.18), y + Inches(0.32), cell_w - Inches(0.18),
                 cell_h - Inches(0.32), clip(hl.desc, 32), size=8, color=BLACK,
                 line_spacing=1.2)

    # ── 投资风险（灰底，2 条横排） ─────────────────────────
    RISK_TOP = Inches(6.18)
    RISK_H = Inches(0.70)
    add_filled_rect(LEFT_X, RISK_TOP, Inches(12.2), RISK_H, GRAY_BG)
    # 标签
    add_text(LEFT_X + Inches(0.1), RISK_TOP, Inches(1.2), RISK_H,
             "投资风险", size=12, bold=True, color=RISK_FG, anchor_top=False)
    risks = payload.risks[:2]
    risk_area_x = LEFT_X + Inches(1.35)
    risk_area_w = Inches(12.2) - Inches(1.35)
    rw = risk_area_w // max(len(risks), 1)
    for i, rk in enumerate(risks):
        x = risk_area_x + rw * i
        box = s.shapes.add_textbox(x, RISK_TOP, rw, RISK_H)
        box.line.fill.background()
        box.fill.background()
        tf = box.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.10)
        p = tf.paragraphs[0]
        p.line_spacing = 1.15
        r1 = p.add_run()
        r1.text = f"{clip(rk.title, 8)}： "
        _set_cn_font(r1, 8, bold=True, color=RISK_FG)
        r2 = p.add_run()
        r2.text = clip(rk.desc, 30)
        _set_cn_font(r2, 8, color=BLACK)

    # ── 页脚小条 ────────────────────────────────────────────
    FOOT_TOP = Inches(7.02)
    f = payload.footer
    foot_text = (
        f"成立年份 {clip(f.founded, 12)}　·　团队规模 {clip(f.team_size, 12)}"
        f"　·　累计融资 {clip(f.funding_total, 14)}　·　{clip(f.ai_grade, 18)}"
    )
    add_text(LEFT_X, FOOT_TOP, Inches(12.2), Inches(0.35),
             foot_text, size=9, color=GRAY, anchor_top=False)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


@app.post("/generate/onepager")
async def generate_onepager(payload: OnePagerPayload):
    try:
        data = _build_onepager(payload)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"一页 PPT 生成失败: {str(e)}")
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=onepager.pptx"},
    )


@app.post("/generate/project_brief")
async def generate_project_brief(payload: dict):
    """
    项目简报 3 页 deck.
    Body: 严格符合 server/services/project_brief/content_schema.json 的对象.
    版式锁在 project_brief_render.py 中.
    """
    import project_brief_render as pb
    buf = io.BytesIO()
    tmp_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
    try:
        pb.render(payload, tmp_path)
        with open(tmp_path, "rb") as f:
            buf.write(f.read())
    except KeyError as e:
        raise HTTPException(status_code=400, detail=f"内容缺失字段: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"项目简报渲染失败: {e}")
    finally:
        try: os.remove(tmp_path)
        except Exception: pass
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=project_brief.pptx"},
    )


@app.post("/generate/investment_deck")
async def generate_investment_deck(payload: dict):
    """
    可变页数投决材料 / 可研报告 / 尽调汇报 deck.
    Body: 严格符合 server/services/investment_deck/content_schema.json 的对象.
    版式锁在 investment_deck_render.py 中.
    """
    import investment_deck_render as deck
    buf = io.BytesIO()
    tmp_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
    try:
        deck.render(payload, tmp_path)
        with open(tmp_path, "rb") as f:
            buf.write(f.read())
    except KeyError as e:
        raise HTTPException(status_code=400, detail=f"内容缺失字段: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"投决材料渲染失败: {e}")
    finally:
        try: os.remove(tmp_path)
        except Exception: pass
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=investment_deck.pptx"},
    )


@app.post("/generate/investment_snapshot")
async def generate_investment_snapshot(payload: dict):
    """
    一页纸投决速览 · 砖红 A4 横版。
    Body: 严格符合 server/services/investment_snapshot/content_schema.json 的对象。
    版式锁在 investment_snapshot_render.py 中，本端点不做任何样式判断。
    """
    import investment_snapshot_render as snap
    buf = io.BytesIO()
    tmp_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
    try:
        snap.render(payload, tmp_path)
        with open(tmp_path, "rb") as f:
            buf.write(f.read())
    except KeyError as e:
        raise HTTPException(status_code=400, detail=f"内容缺失字段: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"投决速览渲染失败: {e}")
    finally:
        try: os.remove(tmp_path)
        except Exception: pass
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=investment_snapshot.pptx"},
    )


@app.get("/health")
async def health():
    return {"status": "ok", "service": "doc-extraction"}
